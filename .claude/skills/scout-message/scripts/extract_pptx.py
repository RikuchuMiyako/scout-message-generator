#!/usr/bin/env python3
"""PPTX からテキストを抽出して標準出力に出す（採用ピッチ資料の読み取り用）。

使い方:
    python scripts/extract_pptx.py <file.pptx>

python-pptx が無い場合は `pip install python-pptx` を実行する。
"""

from __future__ import annotations

import sys


def extract(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            chunks.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(chunks).strip()


def main() -> int:
    if len(sys.argv) != 2:
        print("usage: python scripts/extract_pptx.py <file.pptx>", file=sys.stderr)
        return 2
    try:
        print(extract(sys.argv[1]))
    except ImportError:
        print("python-pptx が必要です: pip install python-pptx", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
