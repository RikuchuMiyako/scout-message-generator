"""PDF からのテキスト抽出ユーティリティ。"""

from __future__ import annotations

import io


def extract_text_from_pdf(data: bytes) -> str:
    """PDF のバイト列からプレーンテキストを抽出して返す。

    求人票・求職者情報の PDF を読み取る用途。pypdf を遅延 import することで、
    PDF を使わない場合は依存を要求しない。
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - 依存未導入時の案内
        raise RuntimeError(
            "PDF を読み取るには pypdf が必要です。`pip install pypdf` を実行してください。"
        ) from exc

    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages).strip()


def extract_text_from_pptx(data: bytes) -> str:
    """PPTX のバイト列からテキストを抽出して返す（採用ピッチ資料など）。"""
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - 依存未導入時の案内
        raise RuntimeError(
            "PPTX を読み取るには python-pptx が必要です。"
            "`pip install python-pptx` を実行してください。"
        ) from exc

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    chunks.append(text)
    return "\n".join(chunks).strip()


def extract_document_text(filename: str, data: bytes) -> str:
    """拡張子に応じて PDF / PPTX / テキストを判定して抽出する。"""
    lower = (filename or "").lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf(data)
    if lower.endswith(".pptx"):
        return extract_text_from_pptx(data)
    # それ以外はプレーンテキストとして解釈
    return data.decode("utf-8", errors="replace").strip()

